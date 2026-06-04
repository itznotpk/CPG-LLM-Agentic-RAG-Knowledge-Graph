"""Generate student-friendly rPPG presentation slides."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── Colour palette ──────────────────────────────────────────
# Store as (r,g,b) tuples so we can do arithmetic, convert with rgb()
def rgb(r, g, b): return RGBColor(r, g, b)
def dim(t, d=4):  return rgb(t[0]//d, t[1]//d, t[2]//d)

_BG_DARK  = (0x0D, 0x11, 0x17)
_ACCENT   = (0x00, 0xFF, 0x88)
_ACCENT2  = (0x38, 0xBD, 0xF8)
_RED      = (0xEF, 0x44, 0x44)
_ORANGE   = (0xFB, 0x92, 0x3A)
_PURPLE   = (0xA7, 0x8B, 0xFA)
_PINK     = (0xF4, 0x72, 0xB6)
_WHITE    = (0xFF, 0xFF, 0xFF)
_GREY     = (0x94, 0xA3, 0xB8)
_YELLOW   = (0xFB, 0xBF, 0x24)

BG_DARK  = rgb(*_BG_DARK)
BG_CARD  = rgb(0x1A, 0x22, 0x2E)
ACCENT   = rgb(*_ACCENT)
ACCENT2  = rgb(*_ACCENT2)
RED      = rgb(*_RED)
ORANGE   = rgb(*_ORANGE)
PURPLE   = rgb(*_PURPLE)
PINK     = rgb(*_PINK)
WHITE    = rgb(*_WHITE)
GREY     = rgb(*_GREY)
YELLOW   = rgb(*_YELLOW)


def blank_slide(prs):
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)


def bg(slide, color=BG_DARK):
    """Fill slide background."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, left, top, width, height, bg_color=None, border_color=None, border_width=Pt(0)):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.line.width = border_width
    if border_color:
        shape.line.color.rgb = border_color
    else:
        shape.line.fill.background()
    if bg_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
    else:
        shape.fill.background()
    return shape


def txt(slide, text, left, top, width, height,
        size=24, bold=False, color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb


def add_para(tf, text, size=20, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, space_before=Pt(4)):
    from pptx.oxml.ns import qn
    from lxml import etree
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return p


def slide_number(slide, n, total=8):
    txt(slide, f"{n} / {total}",
        11.8, 7.1, 1.3, 0.35,
        size=10, color=GREY, align=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════
s = blank_slide(prs)
bg(s)

# Green top bar
box(s, 0, 0, 13.33, 0.08, bg_color=ACCENT)

# Big heart icon area
txt(s, "🫀", 0.6, 1.2, 1.5, 1.5, size=72, align=PP_ALIGN.CENTER)

# Title
txt(s, "rPPG Vital Signs Monitor",
    2.2, 1.3, 10, 1.2, size=44, bold=True, color=WHITE)

# Subtitle
txt(s, "How a webcam and a small sensor can measure your heartbeat",
    2.2, 2.6, 10, 0.8, size=22, color=GREY)

# Tag pills
for i, (label, col) in enumerate([
    ("Heart Rate",      _RED),
    ("SpO₂",           _ACCENT2),
    ("Blood Pressure",  _PINK),
    ("Temperature",     _ORANGE),
    ("Respiratory Rate",_PURPLE),
]):
    box(s, 2.2 + i * 2.15, 3.55, 1.95, 0.42,
        bg_color=dim(col, 4), border_color=rgb(*col), border_width=Pt(1))
    txt(s, label, 2.2 + i * 2.15, 3.55, 1.95, 0.42,
        size=13, color=rgb(*col), align=PP_ALIGN.CENTER)

# Bottom credit
txt(s, "For secondary school students  ·  Live demo included",
    0, 6.8, 13.33, 0.5, size=14, color=GREY, align=PP_ALIGN.CENTER)

box(s, 0, 7.42, 13.33, 0.08, bg_color=ACCENT)
slide_number(s, 1)


# ════════════════════════════════════════════════════════════
# SLIDE 2 — What are Vital Signs?
# ════════════════════════════════════════════════════════════
s = blank_slide(prs)
bg(s)
box(s, 0, 0, 13.33, 0.08, bg_color=ACCENT)

txt(s, "What are Vital Signs?", 0.5, 0.25, 12, 0.7,
    size=34, bold=True, color=WHITE)
txt(s, "These are the basic measurements that tell doctors if your body is working properly.",
    0.5, 1.05, 12, 0.5, size=18, color=GREY)

vitals = [
    ("❤️", "Heart Rate",        "How many times your heart beats per minute",       "60 – 100 BPM",    _RED),
    ("🩸", "Blood Oxygen",      "How much oxygen is carried in your blood",          "95 – 100 %",      _ACCENT2),
    ("💉", "Blood Pressure",    "The force of blood pushing against artery walls",   "< 120 / 80 mmHg", _PINK),
    ("🌡️", "Temperature",       "Your body's internal heat level",                  "36.1 – 37.2 °C",  _ORANGE),
    ("🫁", "Respiratory Rate",  "How many breaths you take per minute",              "12 – 20 br/min",  _PURPLE),
]

for i, (icon, name, desc, normal, color) in enumerate(vitals):
    left = 0.4
    top  = 1.75 + i * 1.0
    box(s, left, top, 12.5, 0.85,
        bg_color=dim(color, 5), border_color=rgb(*color), border_width=Pt(1))
    txt(s, icon,   left + 0.1,  top + 0.05, 0.65, 0.75, size=28)
    txt(s, name,   left + 0.85, top + 0.05, 2.8,  0.4,  size=17, bold=True, color=rgb(*color))
    txt(s, desc,   left + 0.85, top + 0.42, 7.5,  0.38, size=13, color=GREY)
    box(s, left + 9.5, top + 0.15, 2.8, 0.52,
        bg_color=dim(color, 6), border_color=rgb(*color), border_width=Pt(1))
    txt(s, f"Normal: {normal}", left + 9.5, top + 0.15, 2.8, 0.52,
        size=13, color=rgb(*color), align=PP_ALIGN.CENTER)

box(s, 0, 7.42, 13.33, 0.08, bg_color=ACCENT)
slide_number(s, 2)


# ════════════════════════════════════════════════════════════
# SLIDE 3 — What is rPPG?
# ════════════════════════════════════════════════════════════
s = blank_slide(prs)
bg(s)
box(s, 0, 0, 13.33, 0.08, bg_color=ACCENT)

txt(s, "What is rPPG?", 0.5, 0.25, 12, 0.7, size=34, bold=True, color=WHITE)
txt(s, "Remote PhotoPlethysmoGraphy  —  measuring pulse without touching you",
    0.5, 1.0, 12, 0.5, size=18, color=GREY)

# Left — simple analogy box
box(s, 0.4, 1.65, 5.9, 5.2,
    bg_color=RGBColor(0x0F, 0x2A, 0x1F),
    border_color=ACCENT, border_width=Pt(1.5))

txt(s, "💡  The simple idea",
    0.6, 1.8, 5.5, 0.5, size=16, bold=True, color=ACCENT)

steps = [
    "1.  Every heartbeat pumps blood through your face",
    "2.  More blood = slightly redder skin colour",
    "3.  A camera takes ~15 photos per second",
    "4.  A computer detects the tiny colour changes",
    "5.  Count the changes → Heart Rate!",
]
for i, step in enumerate(steps):
    txt(s, step, 0.65, 2.4 + i * 0.76, 5.4, 0.7, size=15, color=WHITE)

txt(s, "Your eye cannot see this change.\nThe camera can! 📷",
    0.65, 6.1, 5.4, 0.6, size=14, color=YELLOW)

# Right — key facts
box(s, 6.8, 1.65, 6.1, 2.4,
    bg_color=RGBColor(0x10, 0x1E, 0x35),
    border_color=ACCENT2, border_width=Pt(1.5))
txt(s, "✅  What rPPG can measure",
    7.0, 1.8, 5.7, 0.5, size=16, bold=True, color=ACCENT2)
for i, item in enumerate(["Heart Rate", "SpO₂ (estimate)", "Respiratory Rate", "Blood Pressure (estimate)"]):
    txt(s, f"   •  {item}", 7.0, 2.4 + i * 0.48, 5.5, 0.45, size=15, color=WHITE)

box(s, 6.8, 4.25, 6.1, 2.6,
    bg_color=RGBColor(0x2A, 0x10, 0x10),
    border_color=RED, border_width=Pt(1.5))
txt(s, "❌  What rPPG cannot replace",
    7.0, 4.4, 5.7, 0.5, size=16, bold=True, color=RED)
for i, item in enumerate(["Clinical-grade blood pressure", "ECG / heart waveform", "Blood glucose level"]):
    txt(s, f"   •  {item}", 7.0, 4.95 + i * 0.48, 5.5, 0.45, size=15, color=GREY)

box(s, 0, 7.42, 13.33, 0.08, bg_color=ACCENT)
slide_number(s, 3)


# ════════════════════════════════════════════════════════════
# SLIDE 4 — Hardware
# ════════════════════════════════════════════════════════════
s = blank_slide(prs)
bg(s)
box(s, 0, 0, 13.33, 0.08, bg_color=ACCENT)

txt(s, "The Hardware", 0.5, 0.25, 12, 0.7, size=34, bold=True, color=WHITE)
txt(s, "Two small devices that work together to read your vital signs",
    0.5, 1.0, 12, 0.5, size=18, color=GREY)

# ESP32 card
box(s, 0.4, 1.65, 5.9, 5.3,
    bg_color=RGBColor(0x10, 0x20, 0x10),
    border_color=ACCENT, border_width=Pt(2))
txt(s, "⚙️  ESP32 DevKit",
    0.6, 1.8, 5.5, 0.55, size=20, bold=True, color=ACCENT)
txt(s, "The brain of the system",
    0.6, 2.4, 5.5, 0.4, size=14, color=GREY)

esp_facts = [
    "A tiny computer chip (like a mini PC)",
    "Connects to WiFi to send data",
    "Reads data from the MAX30100 sensor",
    "Sends vital signs to your browser",
    "Costs about RM 15–25",
]
for i, f in enumerate(esp_facts):
    txt(s, f"  ✦  {f}", 0.6, 2.95 + i * 0.6, 5.5, 0.55, size=14, color=WHITE)

# MAX30100 card
box(s, 6.9, 1.65, 5.9, 5.3,
    bg_color=RGBColor(0x10, 0x10, 0x25),
    border_color=PURPLE, border_width=Pt(2))
txt(s, "💡  MAX30100 Sensor",
    7.1, 1.8, 5.5, 0.55, size=20, bold=True, color=PURPLE)
txt(s, "The finger sensor",
    7.1, 2.4, 5.5, 0.4, size=14, color=GREY)

max_facts = [
    "Shines red + infrared LED onto your finger",
    "Measures how much light bounces back",
    "More blood = less light reflected",
    "Detects Heart Rate and SpO₂",
    "Costs about RM 8–15",
]
for i, f in enumerate(max_facts):
    txt(s, f"  ✦  {f}", 7.1, 2.95 + i * 0.6, 5.5, 0.55, size=14, color=WHITE)

box(s, 0, 7.42, 13.33, 0.08, bg_color=ACCENT)
slide_number(s, 4)


# ════════════════════════════════════════════════════════════
# SLIDE 5 — Wiring
# ════════════════════════════════════════════════════════════
s = blank_slide(prs)
bg(s)
box(s, 0, 0, 13.33, 0.08, bg_color=ACCENT)

txt(s, "How to Connect the Hardware", 0.5, 0.25, 12, 0.7,
    size=34, bold=True, color=WHITE)
txt(s, "Just 4 wires from the MAX30100 to the ESP32",
    0.5, 1.0, 12, 0.5, size=18, color=GREY)

# Wiring table
headers = ["MAX30100 Pin", "→", "ESP32 Pin", "Purpose"]
col_w   = [3.0, 0.6, 3.0, 5.3]
col_x   = [0.4, 3.4, 4.0, 7.1]

# Header row
for j, (h, w, x) in enumerate(zip(headers, col_w, col_x)):
    box(s, x, 1.75, w, 0.5,
        bg_color=RGBColor(0x1E, 0x3A, 0x2E) if j != 1 else BG_DARK,
        border_color=ACCENT if j != 1 else None,
        border_width=Pt(1) if j != 1 else Pt(0))
    txt(s, h, x, 1.75, w, 0.5,
        size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

rows = [
    ("VIN",   "→",  "3.3V",         "Power supply",           YELLOW),
    ("GND",   "→",  "GND",          "Common ground",           GREY),
    ("SDA",   "→",  "GPIO 21",      "Data line (I2C)",         ACCENT2),
    ("SCL",   "→",  "GPIO 22",      "Clock line (I2C)",        PURPLE),
]
row_colors = [
    RGBColor(0x1A, 0x18, 0x05),
    RGBColor(0x10, 0x10, 0x10),
    RGBColor(0x05, 0x18, 0x1A),
    RGBColor(0x15, 0x05, 0x1A),
]
for i, (pin, arrow, esp, purpose, color) in enumerate(rows):
    top = 2.35 + i * 0.72
    for j, (val, w, x) in enumerate(zip([pin, arrow, esp, purpose], col_w, col_x)):
        if j != 1:
            box(s, x, top, w, 0.62,
                bg_color=row_colors[i],
                border_color=color, border_width=Pt(1))
        txt(s, val, x, top, w, 0.62,
            size=16 if j != 1 else 20,
            bold=(j == 0),
            color=color if j != 1 else WHITE,
            align=PP_ALIGN.CENTER)

# Note
box(s, 0.4, 5.5, 12.5, 0.7,
    bg_color=RGBColor(0x1A, 0x16, 0x05),
    border_color=YELLOW, border_width=Pt(1))
txt(s, "⚠️   Use 3.3V — NOT 5V.  The MAX30100 can be damaged by 5V.",
    0.6, 5.55, 12.1, 0.6, size=16, color=YELLOW)

# I2C note
txt(s, "I2C = a communication protocol that lets two devices talk using just 2 wires (SDA + SCL)",
    0.4, 6.35, 12.5, 0.5, size=13, color=GREY)

box(s, 0, 7.42, 13.33, 0.08, bg_color=ACCENT)
slide_number(s, 5)


# ════════════════════════════════════════════════════════════
# SLIDE 6 — System Architecture
# ════════════════════════════════════════════════════════════
s = blank_slide(prs)
bg(s)
box(s, 0, 0, 13.33, 0.08, bg_color=ACCENT)

txt(s, "How the System Works Together", 0.5, 0.25, 12, 0.7,
    size=34, bold=True, color=WHITE)
txt(s, "Three parts talking to each other over WiFi",
    0.5, 1.0, 12, 0.5, size=18, color=GREY)

components = [
    (0.3,  "🖥️",  "Your Webcam",        "Captures your face\n~15 frames/sec",                    ACCENT2,  RGBColor(0x05,0x18,0x22)),
    (4.6,  "⚙️",  "FastAPI Server",     "Python backend\nRuns on your PC",                        ACCENT,   RGBColor(0x05,0x20,0x10)),
    (8.9,  "📱",  "Browser Dashboard", "Shows all vitals\nLive updates",                          PURPLE,   RGBColor(0x15,0x08,0x25)),
]
sensor_box = (0.3, "💉", "ESP32 + MAX30100", "Finger sensor\nSends via WiFi HTTP", YELLOW, RGBColor(0x22,0x18,0x05))

for left, icon, title, desc, color, bgc in components:
    box(s, left, 1.85, 3.9, 2.8, bg_color=bgc, border_color=color, border_width=Pt(2))
    txt(s, icon,  left+0.15, 1.95, 0.8,  0.8,  size=32)
    txt(s, title, left+0.15, 2.85, 3.6,  0.5,  size=16, bold=True, color=color)
    txt(s, desc,  left+0.15, 3.38, 3.6,  0.9,  size=13, color=GREY)

# Arrows between components
for ax in [4.25, 8.55]:
    txt(s, "→", ax, 2.7, 0.5, 0.6, size=28, bold=True, color=GREY, align=PP_ALIGN.CENTER)

# ESP32 box at bottom centre
box(s, 4.6, 5.1, 3.9, 1.8,
    bg_color=RGBColor(0x22, 0x18, 0x05),
    border_color=YELLOW, border_width=Pt(2))
txt(s, "💉", 4.75, 5.2, 0.8, 0.7, size=28)
txt(s, "ESP32 + MAX30100",  4.75, 5.95, 3.6, 0.45, size=15, bold=True, color=YELLOW)
txt(s, "Finger sensor  ·  sends via WiFi HTTP POST", 4.75, 6.42, 3.6, 0.4, size=12, color=GREY)

# Arrow from ESP32 up to server
txt(s, "↑ WiFi", 6.2, 4.7, 1.0, 0.4, size=13, color=YELLOW, align=PP_ALIGN.CENTER)

# Data flow labels
txt(s, "Video frames\n(WebSocket)", 3.5, 2.1, 1.3, 0.8, size=11, color=GREY, align=PP_ALIGN.CENTER)
txt(s, "Vitals data\n(WebSocket)", 7.9, 2.1, 1.2, 0.8, size=11, color=GREY, align=PP_ALIGN.CENTER)

box(s, 0, 7.42, 13.33, 0.08, bg_color=ACCENT)
slide_number(s, 6)


# ════════════════════════════════════════════════════════════
# SLIDE 7 — Live Demo Guide
# ════════════════════════════════════════════════════════════
s = blank_slide(prs)
bg(s)
box(s, 0, 0, 13.33, 0.08, bg_color=ACCENT)

txt(s, "Live Demo — Try It Yourself!", 0.5, 0.25, 12, 0.7,
    size=34, bold=True, color=WHITE)
txt(s, "Follow these steps in order",
    0.5, 1.0, 12, 0.4, size=18, color=GREY)

steps = [
    ("1", "Open the browser",       "Go to  http://172.20.10.3:8090",                   _ACCENT2, rgb(0x05,0x18,0x22)),
    ("2", "Allow camera access",    "Click 'Allow' when the browser asks for camera",   _ACCENT,  rgb(0x05,0x20,0x10)),
    ("3", "Look at the camera",     "Keep your face in the green box on screen",        _YELLOW,  rgb(0x22,0x18,0x05)),
    ("4", "Place finger on sensor", "Gently rest fingertip on the MAX30100",            _PURPLE,  rgb(0x15,0x08,0x25)),
    ("5", "Wait ~10 seconds",       "Numbers will stabilise — that's your heart rate!", _ORANGE,  rgb(0x22,0x10,0x05)),
]

for i, (num, title, desc, color, bgc) in enumerate(steps):
    top = 1.6 + i * 1.1
    box(s, 0.4, top, 12.5, 0.95, bg_color=bgc, border_color=rgb(*color), border_width=Pt(1.5))
    box(s, 0.55, top + 0.18, 0.6, 0.6,
        bg_color=dim(color, 3), border_color=rgb(*color), border_width=Pt(1.5))
    txt(s, num,   0.55, top + 0.18, 0.6,  0.6,  size=18, bold=True, color=rgb(*color), align=PP_ALIGN.CENTER)
    txt(s, title, 1.3,  top + 0.08, 3.5,  0.45, size=17, bold=True, color=rgb(*color))
    txt(s, desc,  1.3,  top + 0.52, 11.1, 0.38, size=14, color=GREY)

box(s, 0, 7.42, 13.33, 0.08, bg_color=ACCENT)
slide_number(s, 7)


# ════════════════════════════════════════════════════════════
# SLIDE 8 — Key Takeaways
# ════════════════════════════════════════════════════════════
s = blank_slide(prs)
bg(s)
box(s, 0, 0, 13.33, 0.08, bg_color=ACCENT)

txt(s, "What Did We Learn?", 0.5, 0.25, 12, 0.7,
    size=34, bold=True, color=WHITE)

takeaways = [
    ("🫀", "rPPG uses a camera to detect your pulse",           "The camera sees tiny colour changes in your skin caused by blood flow",    _ACCENT),
    ("💡", "No contact needed for heart rate",                   "Just look at the camera — technology does the rest",                       _ACCENT2),
    ("⚙️", "ESP32 + MAX30100 adds accuracy",                    "The finger sensor gives a more reliable reading alongside the webcam",     _YELLOW),
    ("📊", "Estimates vs clinical measurements",                "rPPG gives good estimates but is not a replacement for hospital equipment", _ORANGE),
    ("🔮", "The future of health monitoring",                   "Smartwatches, phones and laptops may soon measure vitals automatically",    _PURPLE),
]

for i, (icon, title, desc, color) in enumerate(takeaways):
    top = 1.3 + i * 1.12
    box(s, 0.4, top, 12.5, 1.0,
        bg_color=dim(color, 6), border_color=rgb(*color), border_width=Pt(1))
    txt(s, icon,  0.55, top + 0.1,  0.75, 0.8,  size=28, align=PP_ALIGN.CENTER)
    txt(s, title, 1.45, top + 0.08, 11.2, 0.42, size=17, bold=True, color=rgb(*color))
    txt(s, desc,  1.45, top + 0.52, 11.2, 0.42, size=13, color=GREY)

box(s, 0, 7.42, 13.33, 0.08, bg_color=ACCENT)
slide_number(s, 8)


# ── Save ────────────────────────────────────────────────────
out = r"d:\CPG-LLM-Agentic-RAG-Knowledge-Graph\rppg_poc\rPPG_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
